#!/usr/bin/env python3
"""
Minimal example: McKinsey-style PPT with Cover + Content + Source slides.
Uses the design system defined in SKILL.md (v1.2.0).
"""

import os
import shutil
import subprocess
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color Palette ──
NAVY      = RGBColor(0x05, 0x1C, 0x2C)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY  = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)

# ── Font Sizes ──
TITLE_SIZE      = Pt(22)
BODY_SIZE       = Pt(14)
SUB_HEADER_SIZE = Pt(18)
HEADER_SIZE     = Pt(28)
SMALL_SIZE      = Pt(9)

# ── Helper Functions ──

def _clean_shape(shape):
    """Remove p:style from any shape to prevent effect references."""
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)

def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)

def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = font_name
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for run in p.runs:
        set_ea_font(run, ea_font)
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)  # CRITICAL: remove p:style
    return shape

def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    """Draw a horizontal line using a thin rectangle (no connector)."""
    h = max(int(thickness), Emu(6350))  # minimum ~0.5pt
    return add_rect(slide, x, y, length, h, color)

def add_oval(slide, x, y, letter, size=Inches(0.45),
             bg=NAVY, fg=WHITE):
    """Add a circle label with a letter (e.g. 'A', '1').
    Uses Arial font to match body text consistency."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'KaiTi')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean_shape(c)  # CRITICAL: remove p:style
    return c

def add_action_title(slide, text, title_size=Pt(22)):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, Inches(0.8), Inches(1.05), Inches(11.7),
              color=BLACK, thickness=Pt(0.5))

def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=Pt(9), font_color=MED_GRAY)

def full_cleanup(outpath):
    """Remove ALL p:style from every slide + theme shadows/3D."""
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    # Remove ALL p:style elements from all shapes/connectors
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    # Remove shadows and 3D from theme
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

def deliver_to_channel(outpath, slide_count):
    """Send generated PPTX back to user's chat channel via OpenClaw media pipeline.
    Falls back gracefully if not running in a channel context."""
    if not shutil.which('openclaw'):
        print(f'[deliver] openclaw CLI not found, skipping channel delivery')
        print(f'[deliver] File saved locally: {outpath}')
        return False

    size_kb = os.path.getsize(outpath) / 1024
    caption = f'✅ PPT generated — {slide_count} slides, {size_kb:.0f} KB'

    try:
        result = subprocess.run(
            ['openclaw', 'message', 'send',
             '--media', outpath,
             '--message', caption],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            print(f'[deliver] Sent to channel: {outpath}')
            return True
        else:
            print(f'[deliver] Channel send failed: {result.stderr}')
            print(f'[deliver] File saved locally: {outpath}')
            return False
    except Exception as e:
        print(f'[deliver] Error: {e}')
        print(f'[deliver] File saved locally: {outpath}')
        return False

# ── Build Presentation ──

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, prs.slide_width, Inches(0.05), NAVY)
    add_text(s1, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
             'Sample Presentation', font_size=Pt(44), font_name='Georgia',
             font_color=NAVY, bold=True)
    add_text(s1, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             'McKinsey-style Design System Demo', font_size=Pt(24),
             font_color=DARK_GRAY)
    add_text(s1, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             'Minimal Example  |  2026', font_size=BODY_SIZE, font_color=MED_GRAY)
    add_hline(s1, Inches(1), Inches(6.8), Inches(3), color=NAVY, thickness=Pt(2))

    # Slide 2: Content
    s2 = prs.slides.add_slide(blank)
    add_action_title(s2, 'Key Findings Overview')
    items = [
        'Clean typography hierarchy ensures readability',
        'Flat design with no shadows or 3D effects',
        'Consistent color palette across all slides',
        'Proper East Asian font handling for Chinese text',
    ]
    for i, item in enumerate(items):
        y = Inches(1.6) + Inches(0.6) * i
        add_oval(s2, Inches(0.9), y, str(i + 1))
        add_text(s2, Inches(1.5), y, Inches(10), Inches(0.5), item)
        if i < len(items) - 1:
            add_hline(s2, Inches(0.9), y + Inches(0.55), Inches(11.3), LINE_GRAY)
    add_source(s2, 'Source: Mck-ppt-design-skill v1.2.0')

    # Save & cleanup & deliver
    outpath = 'minimal_output.pptx'
    prs.save(outpath)
    full_cleanup(outpath)
    deliver_to_channel(outpath, len(prs.slides))
    print(f'Created: {outpath} ({os.path.getsize(outpath):,} bytes)')

if __name__ == '__main__':
    main()
