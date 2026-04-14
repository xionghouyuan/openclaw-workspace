#!/usr/bin/env python3
"""银浆印刷加热雨刮卧倒区域方案 — Bain风格 PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Bain palette ──
BAIN_RED   = RGBColor(0xCC, 0x22, 0x29)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x4A, 0x4A, 0x4A)
MED_GRAY   = RGBColor(0x7F, 0x7F, 0x7F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BORDER_G   = RGBColor(0xD9, 0xD9, 0xD9)
CHARCOAL   = RGBColor(0x3D, 0x3D, 0x3D)

# ── Helpers ──
def rb(slide, l, t, w, h, fc=None, oc=None, lw=0.75):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = Pt(lw)
    if fc: shape.fill.solid(); shape.fill.fore_color.rgb = fc
    else: shape.fill.background()
    if oc: shape.line.color.rgb = oc
    else: shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h, size=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, font='Arial'):
    from pptx.util import Pt
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape

def add_red_bar(slide):
    rb(slide, 0, 0, 10, 0.055, fc=BAIN_RED)

def add_footer(slide, src, pg, total):
    txb(slide, f'Source: {src}', 0.4, 7.35, 5, 0.3, size=7.5, color=MED_GRAY)
    txb(slide, f'{pg} / {total}', 9.2, 7.35, 0.6, 0.3, size=8, color=MED_GRAY,
        align=PP_ALIGN.RIGHT)

def add_section_tag(slide, tag, l=0.4, t=0.15):
    txb(slide, tag, l, t, 2, 0.3, size=8, color=MED_GRAY)

# ── Create presentation ──
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)

TOTAL = 5

# ════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════
def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_red_bar(slide)

    # Title block
    txb(slide, '银浆印刷加热雨刮',
        0.5, 1.4, 9, 0.9, size=38, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    txb(slide, '卧倒区域方案',
        0.5, 2.25, 9, 0.9, size=38, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    # Divider
    rb(slide, 0.5, 3.25, 2.5, 0.04, fc=BAIN_RED)

    # Subtitle
    txb(slide, '开发费用周期及计划  /  Development Cost Cycle & Plan',
        0.5, 3.45, 8, 0.5, size=15, color=DARK_GRAY)

    # Key stats row
    stats = [
        ('¥5~10万', '总开发费用'),
        ('-45天',   '进度偏差'),
        ('11/15',   'OTS里程碑'),
        ('3项',     '待决策'),
    ]
    x = 0.5
    for val, lbl in stats:
        rb(slide, x, 4.35, 2.1, 1.3, fc=LIGHT_GRAY)
        txb(slide, val, x+0.15, 4.5, 1.9, 0.6, size=28, bold=True, color=BAIN_RED)
        txb(slide, lbl, x+0.15, 5.1, 1.9, 0.4, size=11, color=DARK_GRAY)
        x += 2.3

    # Date
    txb(slide, '2026.04.13', 0.5, 6.5, 2, 0.3, size=10, color=MED_GRAY)
    add_footer(slide, '银浆印刷加热雨刮卧倒区域方案', 1, TOTAL)

slide_cover()

# ════════════════════════════════════════════════
# SLIDE 2 — PROJECT OVERVIEW
# ════════════════════════════════════════════════
def slide_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(slide)
    add_section_tag(slide, '01 / OVERVIEW')

    # Action title
    txb(slide, 'OTS里程碑 11/15，当前进度滞后 45 天，三项决策待确认',
        0.4, 0.38, 9.2, 0.55, size=17, bold=True, color=BLACK)

    # Divider
    rb(slide, 0.4, 0.97, 9.2, 0.015, fc=BORDER_G)

    # Left column — Cost table
    rb(slide, 0.4, 1.1, 4.4, 0.35, fc=BAIN_RED)
    txb(slide, '开发费用 Development Cost', 0.55, 1.13, 4, 0.3, size=10, bold=True, color=WHITE)

    rows = [
        ('模具单侧 Left or Right',  '¥20,000',  LIGHT_GRAY),
        ('模具双侧 Left + Right',   '¥40,000',  WHITE),
        ('试验单侧 Cabin',           '¥30,000',  LIGHT_GRAY),
        ('试验双侧 Cabin',           '¥60,000',  WHITE),
    ]
    y = 1.47
    for lbl, val, bg in rows:
        rb(slide, 0.4, y, 4.4, 0.42, fc=bg, oc=BORDER_G, lw=0.5)
        txb(slide, lbl, 0.55, y+0.04, 2.8, 0.34, size=10.5, color=DARK_GRAY)
        txb(slide, val, 3.5,  y+0.04, 1.2, 0.34, size=11, bold=True,
            color=BAIN_RED if '¥40' in val or '¥60' in val else BLACK,
            align=PP_ALIGN.RIGHT)
        y += 0.44

    # Source
    txb(slide, '含模具 + 试验费用，详见技术规格', 0.4, 3.38, 4.4, 0.25, size=7.5, color=MED_GRAY)

    # Right top — OTS warning
    rb(slide, 5.1, 1.1, 4.5, 1.55, fc=LIGHT_GRAY, oc=BORDER_G)
    rb(slide, 5.1, 1.1, 0.06, 1.55, fc=BAIN_RED)   # red left bar
    txb(slide, 'OTS 里程碑  /  OTS Milestone', 5.25, 1.15, 4.2, 0.3, size=9.5, bold=True, color=BAIN_RED)
    txb(slide, '11/15', 5.25, 1.45, 3.5, 0.75, size=42, bold=True, color=BAIN_RED)
    txb(slide, '当前滞后 45 天，需紧急压缩周期', 5.25, 2.25, 4.2, 0.3, size=10, color=DARK_GRAY)

    # Right bottom — Decision items
    rb(slide, 5.1, 2.8, 4.5, 1.8, fc=WHITE, oc=BORDER_G)
    rb(slide, 5.1, 2.8, 0.06, 1.8, fc=BAIN_RED)
    txb(slide, '待决策事项  /  Decision Items', 5.25, 2.85, 4.2, 0.3, size=9.5, bold=True, color=BAIN_RED)
    decisions = [
        '① 尽快正式提供雨刮区域加热配置的技术要求',
        '② 确认该配置对应的具体销售区域和车型',
        '③ 确认是否需要开发左/右舵 Cabin 双侧配置',
    ]
    dy = 3.18
    for d in decisions:
        txb(slide, d, 5.25, dy, 4.2, 0.4, size=10, color=DARK_GRAY)
        dy += 0.42

    txb(slide, '需尽快确认，避免影响后续开发节点', 5.25, 4.38, 4.2, 0.25, size=7.5, color=MED_GRAY)

    add_footer(slide, '银浆印刷加热雨刮卧倒区域方案 — 开发费用及里程碑', 2, TOTAL)

slide_overview()

# ════════════════════════════════════════════════
# SLIDE 3 — GANTT CHART
# ════════════════════════════════════════════════
def slide_gantt():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(slide)
    add_section_tag(slide, '02 / TIMELINE')

    txb(slide, '开发周期 M10 ~ M10 SOP，当前进度滞后 45 天风险突出',
        0.4, 0.38, 9.2, 0.55, size=17, bold=True, color=BLACK)
    rb(slide, 0.4, 0.97, 9.2, 0.015, fc=BORDER_G)

    # Phase legend (top)
    phases = [
        ('M10\nKO',      BAIN_RED),
        ('M11\nB',      RGBColor(0x7B, 0x3A, 0x9B)),
        ('M1-3\n过程开发', RGBColor(0x8B, 0x5E, 0x3C)),
        ('M4-6\nS1',    RGBColor(0xC0, 0x5C, 0x9C)),
        ('M7-8\nS2',    RGBColor(0xCC, 0x44, 0x4A)),
        ('M9\nP1',      RGBColor(0x1A, 0x7A, 0x3A)),
        ('M10\nSOP',    RGBColor(0x0A, 0x50, 0x30)),
    ]
    px = 0.4
    pw = 1.27
    for lbl, col in phases:
        rb(slide, px, 1.08, pw-0.06, 0.65, fc=col)
        lines = lbl.split('\n')
        txb(slide, lines[0], px, 1.08+0.05, pw-0.06, 0.3, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if len(lines) > 1:
            txb(slide, lines[1], px, 1.38, pw-0.06, 0.3, size=8.5, color=WHITE, align=PP_ALIGN.CENTER)
        px += pw

    # Month labels
    months = ['M10', 'M11', 'M1', 'M2', 'M3', 'M4', 'M5', 'M6', 'M7', 'M8', 'M9', 'M10']
    mx = 0.4
    mw = 0.75
    for m in months:
        txb(slide, m, mx, 1.78, mw, 0.25, size=8, color=MED_GRAY, align=PP_ALIGN.CENTER)
        mx += mw

    # Vertical grid lines
    for i in range(1, 12):
        vx = 0.4 + i * 0.75
        rb(slide, vx, 1.78, 0.008, 0.015, fc=BORDER_G)

    # Task rows
    tasks = [
        ('数据设计',   0.20, BAIN_RED,   'M4~M5'),
        ('设变流程',   0.25, BAIN_RED,   'M5~M6'),
        ('商务议价',   0.30, BAIN_RED,   'M6~M7'),
        ('零件开发',   0.45, CHARCOAL,   'M7~M9'),
        ('型式试验',   0.55, CHARCOAL,   'M9~M11'),
        ('OTS认可',    0.72, BAIN_RED,   'M11~P+'),
    ]

    ty = 2.1
    for name, start_pct, col, dates in tasks:
        # Task name
        rb(slide, 0.4, ty, 1.2, 0.52, fc=LIGHT_GRAY, oc=BORDER_G)
        txb(slide, name, 0.45, ty+0.06, 1.1, 0.4, size=11, color=BLACK)
        # Bar
        bar_l = 1.65 + start_pct * 7.65
        bar_w = 0.75 * (1 - start_pct) * 1.2  # approximate width
        # Color based on phase
        bar_col = BAIN_RED if name in ('数据设计','设变流程','OTS认可') else CHARCOAL
        rb(slide, bar_l, ty+0.06, 7.3 - bar_l + 1.65, 0.4, fc=bar_col)
        txb(slide, dates, bar_l + 0.08, ty+0.08, 1.5, 0.35, size=8.5, color=WHITE)
        # Row bg
        rb(slide, 1.65, ty, 7.65, 0.52, fc=WHITE, oc=BORDER_G, lw=0.3)
        ty += 0.64

    # Progress label
    rb(slide, 0.4, 6.05, 9.2, 0.015, fc=BORDER_G)
    txb(slide, '当前滞后 45 天，关键路径为型式试验（50天）+ OTS认可（30天）',
        0.4, 6.1, 8, 0.3, size=9.5, color=DARK_GRAY)
    add_footer(slide, '银浆印刷加热雨刮卧倒区域方案 — 开发周期', 3, TOTAL)

slide_gantt()

# ════════════════════════════════════════════════
# SLIDE 4 — RISK & MEASURES
# ════════════════════════════════════════════════
def slide_risk():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(slide)
    add_section_tag(slide, '03 / RISK & MITIGATION')

    txb(slide, '当前开发周期无法满足 SOP，OTS 里程碑 11/15 滞后 45 天',
        0.4, 0.38, 9.2, 0.55, size=17, bold=True, color=BLACK)
    rb(slide, 0.4, 0.97, 9.2, 0.015, fc=BORDER_G)

    # Left: Risk
    rb(slide, 0.4, 1.1, 4.4, 2.2, fc=LIGHT_GRAY, oc=BORDER_G)
    rb(slide, 0.4, 1.1, 0.07, 2.2, fc=BAIN_RED)
    txb(slide, '关键风险  Key Risk', 0.6, 1.15, 4, 0.32, size=10, bold=True, color=BAIN_RED)
    txb(slide, '-45天', 0.6, 1.48, 3, 0.8, size=44, bold=True, color=BAIN_RED)
    txb(slide, '当前进度滞后', 0.6, 2.35, 2, 0.35, size=12, color=DARK_GRAY)
    risks = [
        '• 整体周期无法满足 SOP 要求',
        '• 型式试验（50天）和 OTS 认可（30天）',
        '  为关键路径，无法压缩',
        '• 当前已滞后 45 天，需紧急干预',
    ]
    ry = 2.68
    for r in risks:
        txb(slide, r, 0.6, ry, 4, 0.3, size=9.5, color=DARK_GRAY)
        ry += 0.3

    # Right: Measures
    rb(slide, 5.1, 1.1, 4.5, 4.2, fc=WHITE, oc=BORDER_G)
    txb(slide, '应对措施  Mitigation Measures', 5.25, 1.15, 4.2, 0.32, size=10, bold=True, color=BLACK)

    measures = [
        ('措施 1', '目标压缩 30~45 天',
         '提前启动设计工作，压缩设计周期；提前发起紧急设变流程；推动商务议价；提前排期试验。'),
        ('措施 2', '提前 15 天达到量产',
         '关键型式试验完成后，采用过渡 OTS 认可流程，可提前 15 天达到量产条件。'),
    ]
    my = 1.55
    for num, title, desc in measures:
        rb(slide, 5.1, my, 0.08, 0.7, fc=BAIN_RED)
        txb(slide, num, 5.3, my, 1, 0.28, size=9, bold=True, color=BAIN_RED)
        txb(slide, title, 5.3, my+0.28, 4.1, 0.32, size=12, bold=True, color=BLACK)
        txb(slide, desc, 5.3, my+0.62, 4.1, 0.65, size=9.5, color=DARK_GRAY)
        my += 1.35

    # Bottom key message
    rb(slide, 0.4, 5.5, 9.2, 0.8, fc=LIGHT_GRAY)
    txb(slide, '执行要点：提前启动设计、紧急设变、提前排期试验是压缩周期的关键',
        0.55, 5.58, 8.9, 0.6, size=11, color=DARK_GRAY)

    add_footer(slide, '银浆印刷加热雨刮卧倒区域方案 — 风险与应对', 4, TOTAL)

slide_risk()

# ════════════════════════════════════════════════
# SLIDE 5 — DECISION ITEMS
# ════════════════════════════════════════════════
def slide_decisions():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(slide)
    add_section_tag(slide, '04 / DECISION REQUIRED')

    txb(slide, '三项决策需尽快确认，影响后续开发节点和商务议价',
        0.4, 0.38, 9.2, 0.55, size=17, bold=True, color=BLACK)
    rb(slide, 0.4, 0.97, 9.2, 0.015, fc=BORDER_G)

    items = [
        ('01', '技术要求确认',
         '尽快正式提供雨刮区域加热配置的技术要求，包括功率、温度、位置等具体参数，为设计开发提供依据。',
         BAIN_RED),
        ('02', '销售区域与车型确认',
         '确认该配置对应的具体销售区域和车型，明确目标市场，以便进行针对性的开发和认证。',
         CHARCOAL),
        ('03', '左/右舵 Cabin 双侧配置',
         '确认是否需要开发左舵（Left）和右舵（Right）双侧配置，还是仅需单侧 Cabin 配置。',
         RGBColor(0x8C, 0x8C, 0x8C)),
    ]

    y = 1.12
    for num, title, desc, col in items:
        # Number box
        rb(slide, 0.4, y, 0.65, 1.4, fc=col)
        txb(slide, num, 0.4, y+0.45, 0.65, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Content
        rb(slide, 1.12, y, 8.5, 1.4, fc=LIGHT_GRAY, oc=BORDER_G, lw=0.5)
        txb(slide, title, 1.25, y+0.1, 8.2, 0.42, size=14, bold=True, color=BLACK)
        txb(slide, desc, 1.25, y+0.55, 8.2, 0.75, size=10.5, color=DARK_GRAY)
        y += 1.58

    # Bottom
    rb(slide, 0.4, 5.95, 9.2, 0.6, fc=WHITE, oc=BAIN_RED, lw=1.5)
    txb(slide, '请确认以上三项决策，以便尽快推进开发工作，避免影响 SOP 节点',
        0.55, 6.0, 8.9, 0.5, size=11, bold=True, color=BAIN_RED)

    add_footer(slide, '银浆印刷加热雨刮卧倒区域方案 — 待决策事项', 5, TOTAL)

slide_decisions()

# ── Save ──
out = '/home/xionghouyuan2/银浆印刷加热雨刮卧倒区域方案_Bain.pptx'
prs.save(out)
print(f'Saved: {out}')
