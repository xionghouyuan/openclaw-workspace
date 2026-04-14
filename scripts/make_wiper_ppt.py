#!/usr/bin/env python3
"""
银浆印刷加热雨刮卧倒区域方案 — PPT
Linear-style design: 浅色背景 + 蓝白卡片 + indigo品牌色
"""

import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Linear Colors ──
BG         = RGBColor(0xF7, 0xF8, 0xF8)
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL      = RGBColor(0xF3, 0xF4, 0xF5)
BORDER     = RGBColor(0xE6, 0xE6, 0xE6)
TEXT1      = RGBColor(0x1A, 0x1A, 0x1E)
TEXT2      = RGBColor(0x3C, 0x3C, 0x43)
TEXT3      = RGBColor(0x62, 0x66, 0x6D)
TEXT4      = RGBColor(0x8A, 0x8F, 0x98)
BRAND      = RGBColor(0x5E, 0x6A, 0xD2)
ACCENT     = RGBColor(0x71, 0x70, 0xFF)
GREEN      = RGBColor(0x27, 0xA6, 0x44)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
AMBER      = RGBColor(0xF5, 0xA6, 0x23)
RED        = RGBColor(0xE5, 0x48, 0x4D)
BROWN      = RGBColor(0x8B, 0x5E, 0x3C)
PINK       = RGBColor(0xD4, 0x60, 0x9A)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
SOPG       = RGBColor(0x1A, 0x7A, 0x3A)

LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_RED   = RGBColor(0xFF, 0xEB, 0xEE)
LIGHT_AMBER = RGBColor(0xFF, 0xF3, 0xE0)

def _clean(shape):
    sp = shape._element
    s = sp.find(qn('p:style'))
    if s is not None: sp.remove(s)

def set_ea(run, face='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', face)

def txt(slide, l, t, w, h, text, sz=Pt(12), name='Arial', color=TEXT2,
        bold=False, align=PP_ALIGN.LEFT, ea='KaiTi', anchor=MSO_ANCHOR.TOP,
        halign=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 't')
    for a in ['lIns','tIns','rIns','bIns']: bodyPr.set(a, '45720')
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = sz
    p.font.name = name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_before = Pt(0); p.space_after = Pt(0)
    for r in p.runs: set_ea(r, ea)
    return tb

def rect(slide, l, t, w, h, fc, line=None, radius=False):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shape.adjustments[0] = 0.08
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fc
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    _clean(shape)
    return shape

def oval(slide, x, y, label, sz=Inches(0.32), bg=BRAND, fg=RGBColor(255,255,255), fsz=Pt(11)):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    c.fill.solid(); c.fill.fore_color.rgb = bg; c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = str(label)
    tf.paragraphs[0].font.size = fsz
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for r in tf.paragraphs[0].runs: set_ea(r, 'Arial')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean(c); return c

def hline(slide, x, y, len_, color=BORDER, thick=Pt(0.5)):
    h = max(int(thick), Emu(6350))
    return rect(slide, x, y, len_, h, color)

def src(slide, text, y=Inches(7.1)):
    txt(slide, Inches(0.8), y, Inches(11), Inches(0.3), text, sz=Pt(9), color=TEXT4)

def banner(slide, text, lbl_color=BRAND):
    rect(slide, 0, 0, Inches(13.333), Inches(0.06), lbl_color)

def section_title(slide, lbl, title):
    txt(slide, Inches(0.8), Inches(0.18), Inches(4), Inches(0.35),
        lbl, sz=Pt(10), color=TEXT4)
    txt(slide, Inches(0.8), Inches(0.48), Inches(11), Inches(0.6),
        title, sz=Pt(22), name='Georgia', color=TEXT1, bold=True, ea='KaiTi')
    rect(slide, Inches(0.8), Inches(1.05), Inches(11.5), Pt(0.5), BORDER)

def card(slide, l, t, w, h, bg=SURFACE, line=BORDER):
    rect(slide, l, t, w, h, bg, line)

# ── Slides ──

def s_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    banner(s, '银浆印刷加热雨刮卧倒区域方案', BRAND)
    # Background accent
    rect(s, 0, Inches(6.5), prs.slide_width, Inches(1.0), BRAND)
    # Title
    txt(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
        '银浆印刷加热雨刮卧倒区域方案', sz=Pt(40), name='Georgia', color=TEXT1, bold=True, ea='KaiTi')
    txt(s, Inches(1), Inches(3.2), Inches(11), Inches(0.7),
        '开发费用周期及计划', sz=Pt(26), color=TEXT2, ea='KaiTi')
    rect(s, Inches(1), Inches(4.1), Inches(3.5), Pt(3), BRAND)
    txt(s, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
        'xionghouyuan2  |  2026.04.13', sz=Pt(14), color=TEXT4)
    txt(s, Inches(1), Inches(6.7), Inches(11), Inches(0.5),
        'Linear Style  ·  Project Overview', sz=Pt(14), color=RGBColor(200,200,220), align=PP_ALIGN.CENTER)

def s_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    banner(s, '01 / 概览')
    section_title(s, '01 / OVERVIEW', '项目概览 Project Overview')

    # 4 stat cards
    stats = [
        ('¥5~10万', '总开发费用', BRAND, LIGHT_BLUE),
        ('6', '项目任务数', GREEN, LIGHT_GREEN),
        ('-45天', '进度偏差', RED, LIGHT_RED),
        ('3项', '待决策事项', AMBER, LIGHT_AMBER),
    ]
    for i, (num, lbl, color, light) in enumerate(stats):
        x = Inches(0.8) + Inches(3.0) * i
        card(s, x, Inches(1.4), Inches(2.8), Inches(1.5), SURFACE, BORDER)
        rect(s, x, Inches(1.4), Inches(2.8), Inches(0.06), color)
        txt(s, x+Inches(0.15), Inches(1.55), Inches(2.5), Inches(0.7),
            num, sz=Pt(28), name='Georgia', color=color, bold=True)
        txt(s, x+Inches(0.15), Inches(2.4), Inches(2.5), Inches(0.4),
            lbl, sz=Pt(12), color=TEXT3)

    # Risk box
    card(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.6), LIGHT_RED, BORDER)
    rect(s, Inches(0.8), Inches(3.1), Inches(0.06), Inches(1.6), RED)
    txt(s, Inches(1.05), Inches(3.25), Inches(11), Inches(0.4),
        '⚠  关键风险 Key Risk', sz=Pt(14), color=RED, bold=True, ea='KaiTi')
    txt(s, Inches(1.05), Inches(3.75), Inches(11), Inches(0.9),
        '当前开发周期无法满足 SOP 要求，OTS 里程碑 11/15 滞后 45 天。'
        '需提前启动设计、紧急设变、提前排期试验，目标压缩 30~45 天。',
        sz=Pt(13), color=TEXT2, ea='KaiTi')

    # Decision items
    txt(s, Inches(0.8), Inches(4.9), Inches(11), Inches(0.5),
        '决策项 Decision Items', sz=Pt(16), name='Georgia', color=TEXT1, bold=True, ea='KaiTi')
    items = [
        '尽快正式提供雨刮区域加热配置的技术要求',
        '确认该配置对应的具体销售区域和车型',
        '确认是否需要开发左/右舵 Cabin 双侧配置',
    ]
    for i, item in enumerate(items):
        y = Inches(5.45) + Inches(0.48) * i
        oval(s, Inches(0.8), y, str(i+1), bg=BRAND)
        txt(s, Inches(1.3), y, Inches(10.5), Inches(0.4), item, sz=Pt(13), color=TEXT2, ea='KaiTi')
        if i < len(items)-1:
            hline(s, Inches(0.8), y+Inches(0.43), Inches(11.5), BORDER)

    src(s, '银浆印刷加热雨刮卧倒区域方案')

def s_cost(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    banner(s, '02 / 费用')
    section_title(s, '02 / COST', '开发费用 Development Cost')

    # Left: mold cost
    card(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8), SURFACE, BORDER)
    txt(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
        '🛠  模具/夹具费用', sz=Pt(14), color=TEXT1, bold=True, ea='KaiTi')
    hline(s, Inches(0.9), Inches(1.9), Inches(5.2), BORDER)
    rows1 = [('单侧 Cabin（Left or Right）', '¥20,000', TEXT1),
             ('双侧 Cabin（Left + Right）', '¥40,000', BRAND)]
    for i, (lbl, val, col) in enumerate(rows1):
        y = Inches(2.05) + Inches(0.65) * i
        txt(s, Inches(1.0), y, Inches(3.2), Inches(0.4), lbl, sz=Pt(12), color=TEXT3, ea='KaiTi')
        txt(s, Inches(4.2), y, Inches(2), Inches(0.4), val, sz=Pt(16), color=col, bold=True)
        if i < len(rows1)-1:
            hline(s, Inches(1.0), y+Inches(0.55), Inches(5.2), BORDER)

    # Right: test cost
    card(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.8), SURFACE, BORDER)
    txt(s, Inches(7.0), Inches(1.4), Inches(5), Inches(0.5),
        '🔬  试验费用', sz=Pt(14), color=TEXT1, bold=True, ea='KaiTi')
    hline(s, Inches(6.9), Inches(1.9), Inches(5.2), BORDER)
    rows2 = [('单侧 Cabin（Left or Right）', '¥30,000', TEXT1),
             ('双侧 Cabin（Left + Right）', '¥60,000', BRAND)]
    for i, (lbl, val, col) in enumerate(rows2):
        y = Inches(2.05) + Inches(0.65) * i
        txt(s, Inches(7.0), y, Inches(3.2), Inches(0.4), lbl, sz=Pt(12), color=TEXT3, ea='KaiTi')
        txt(s, Inches(10.2), y, Inches(2), Inches(0.4), val, sz=Pt(16), color=col, bold=True)
        if i < len(rows2)-1:
            hline(s, Inches(7.0), y+Inches(0.55), Inches(5.2), BORDER)

    # OTS highlight
    card(s, Inches(0.8), Inches(4.3), Inches(5.5), Inches(1.9), LIGHT_RED, BORDER)
    rect(s, Inches(0.8), Inches(4.3), Inches(0.06), Inches(1.9), RED)
    txt(s, Inches(1.05), Inches(4.4), Inches(5), Inches(0.4),
        '⚠  OTS 里程碑（计划）', sz=Pt(12), color=RED, bold=True, ea='KaiTi')
    txt(s, Inches(1.05), Inches(4.9), Inches(5), Inches(0.8),
        '11/15', sz=Pt(44), name='Georgia', color=RED, bold=True)
    txt(s, Inches(1.05), Inches(5.75), Inches(5), Inches(0.4),
        '当前滞后 45 天，急需压缩周期', sz=Pt(12), color=TEXT3, ea='KaiTi')

    # Summary
    card(s, Inches(6.8), Inches(4.3), Inches(5.5), Inches(1.9), PANEL, BORDER)
    txt(s, Inches(7.0), Inches(4.45), Inches(5), Inches(0.4),
        '费用汇总 Summary', sz=Pt(14), color=TEXT1, bold=True, ea='KaiTi')
    hline(s, Inches(6.9), Inches(4.9), Inches(5.2), BORDER)
    rows3 = [('模具+试验（单侧）', '¥50,000'), ('模具+试验（双侧）', '¥100,000')]
    for i, (lbl, val) in enumerate(rows3):
        y = Inches(5.05) + Inches(0.5) * i
        txt(s, Inches(7.0), y, Inches(3), Inches(0.4), lbl, sz=Pt(12), color=TEXT3, ea='KaiTi')
        txt(s, Inches(10.2), y, Inches(2), Inches(0.4), val, sz=Pt(14), color=BRAND, bold=True)
        if i < len(rows3)-1:
            hline(s, Inches(7.0), y+Inches(0.42), Inches(5.2), BORDER)

    src(s, '银浆印刷加热雨刮卧倒区域方案')

def s_gantt(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    banner(s, '03 / 时间线')
    section_title(s, '03 / TIMELINE', '开发周期甘特图 Development Gantt')

    # Phase legend row
    phases = [
        ('M10 KO', AMBER), ('M11 B', PURPLE), ('M1-3 过程', BROWN),
        ('M4-6 S1', PINK), ('M7-8 S2', RED), ('M9 P1', GREEN), ('M10 SOP', SOPG),
    ]
    total_w = Inches(10.0)
    pw = total_w / len(phases)
    for i, (lbl, color) in enumerate(phases):
        x = Inches(1.5) + pw * i
        rect(s, x, Inches(1.25), pw - Inches(0.02), Inches(0.38), color)
        txt(s, x, Inches(1.27), pw, Inches(0.35),
            lbl, sz=Pt(9), color=SURFACE, bold=True, align=PP_ALIGN.CENTER, ea='KaiTi')

    # Month labels
    months = ['M10','M11','M1','M4','M5','M6','M7','M8','M9','M10']
    for i, m in enumerate(months):
        x = Inches(1.5) + total_w/len(months)*i
        bold = m == 'M4'
        col = BRAND if m == 'M4' else TEXT4
        txt(s, x, Inches(1.7), total_w/len(months), Inches(0.3),
            m, sz=Pt(9), color=col, bold=bold, align=PP_ALIGN.CENTER)

    # Vertical grid lines
    for i in range(11):
        x = Inches(1.5) + total_w/10*i
        rect(s, x, Inches(2.0), Pt(0.5), Inches(4.6), BORDER)

    # Task bars
    tasks = [
        ('数据设计',    20, '4/01~5/01',  RED,   0.30, 0.16),
        ('设变流程',    20, '5/01~6/01',  RED,   0.38, 0.16),
        ('商务议价',    25, '6/01~7/01',  RED,   0.46, 0.13),
        ('零件开发',    50, '7/01~9/01',  RED,   0.54, 0.26),
        ('型式试验',    50, '9/01~11/01', RED,   0.70, 0.26),
        ('OTS认可',     30, '11/01~11/30', BRAND, 0.85, 0.13),
    ]
    for i, (name, days, dates, color, left, width) in enumerate(tasks):
        y = Inches(2.1) + Inches(0.72)*i
        txt(s, Inches(0.2), y, Inches(1.3), Inches(0.6),
            name, sz=Pt(12), color=TEXT2, ea='KaiTi')
        bx = Inches(1.5) + total_w*left
        bw = total_w*width
        rect(s, bx, y+Inches(0.07), bw, Inches(0.32), color)
        txt(s, bx+Inches(0.08), y+Inches(0.1), bw-Inches(0.1), Inches(0.28),
            f'{days}天  {dates}', sz=Pt(9), color=SURFACE, bold=True, ea='KaiTi')
        hline(s, Inches(0.2), y+Inches(0.65), Inches(12.5), BORDER)

    # Current line
    cx = Inches(1.5) + total_w*0.34
    rect(s, cx, Inches(2.0), Pt(1.5), Inches(4.6), RED)
    txt(s, cx-Inches(0.2), Inches(6.6), Inches(1), Inches(0.3),
        '当前', sz=Pt(9), color=RED, bold=True, align=PP_ALIGN.CENTER)

    src(s, '银浆印刷加热雨刮卧倒区域方案')

def s_measures(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    banner(s, '04 / 措施')
    section_title(s, '04 / MEASURES', '应对措施 Mitigation')

    measures = [
        ('1', '目标压缩 30~45 天', AMBER, LIGHT_AMBER,
         '提前启动设计工作，压缩设计周期；提前发起紧急设变流程；推动商务议价；提前排期试验。'),
        ('2', '提前 15 天达到量产', EMERALD, LIGHT_GREEN,
         '关键型式试验完成后，采用过渡 OTS 认可流程，可提前 15 天达到量产条件。'),
    ]
    for i, (num, title, color, light, desc) in enumerate(measures):
        x = Inches(0.8) + Inches(6.15)*i
        card(s, x, Inches(1.3), Inches(5.7), Inches(2.6), SURFACE, BORDER)
        rect(s, x, Inches(1.3), Inches(5.7), Inches(0.06), color)
        oval(s, x+Inches(0.2), Inches(1.5), num, bg=color, sz=Inches(0.42), fsz=Pt(15))
        txt(s, x+Inches(0.75), Inches(1.52), Inches(4.7), Inches(0.5),
            title, sz=Pt(16), color=color, bold=True, ea='KaiTi')
        txt(s, x+Inches(0.2), Inches(2.15), Inches(5.3), Inches(1.6),
            desc, sz=Pt(13), color=TEXT2, ea='KaiTi')

    # Bottom stats
    stats = [
        ('¥5~10万', '总开发费用', BRAND),
        ('M10', '目标 SOP', EMERALD),
        ('-45天', '当前偏差', RED),
        ('30~45天', '目标压缩', AMBER),
    ]
    for i, (num, lbl, color) in enumerate(stats):
        x = Inches(0.8) + Inches(3.0)*i
        card(s, x, Inches(4.2), Inches(2.8), Inches(1.4), PANEL, BORDER)
        txt(s, x, Inches(4.3), Inches(2.8), Inches(0.7),
            num, sz=Pt(26), name='Georgia', color=color, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x, Inches(5.0), Inches(2.8), Inches(0.4),
            lbl, sz=Pt(12), color=TEXT3, align=PP_ALIGN.CENTER)

    # Decision
    card(s, Inches(0.8), Inches(5.85), Inches(11.5), Inches(1.15), LIGHT_BLUE, BORDER)
    rect(s, Inches(0.8), Inches(5.85), Inches(0.06), Inches(1.15), BRAND)
    txt(s, Inches(1.05), Inches(5.9), Inches(11), Inches(0.4),
        '决策项 Decision Items', sz=Pt(13), color=BRAND, bold=True, ea='KaiTi')
    items_str = '① 尽快正式提供雨刮区域加热配置的技术要求   ② 确认具体销售区域和车型   ③ 确认是否需要双侧 Cabin 配置'
    txt(s, Inches(1.05), Inches(6.35), Inches(11), Inches(0.5),
        items_str, sz=Pt(12), color=TEXT2, ea='KaiTi')

    src(s, 'xionghouyuan2  |  2026.04.13')

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s_cover(prs)
    s_overview(prs)
    s_cost(prs)
    s_gantt(prs)
    s_measures(prs)

    outpath = '/home/xionghouyuan2/workplan_project.pptx'
    prs.save(outpath)

    # Cleanup
    tmp = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    nsp = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    nsa = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for st in root.findall(f'.//{{{nsp}}}style'):
                        st.getparent().remove(st)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw','innerShdw','scene3d','sp3d']:
                            for el in root.findall(f'.//{{{nsa}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmp, outpath)
    print(f'Created: {outpath} ({os.path.getsize(outpath):,} bytes)')

if __name__ == '__main__':
    main()
