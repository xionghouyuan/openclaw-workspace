#!/usr/bin/env python3
"""Embed Linear-style HTML screenshots into PPTX"""
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

slides = [
    ('/tmp/slide1.jpg', '封面 - 项目概览'),
    ('/tmp/slide2.jpg', '甘特图 + 风险与决策'),
]

for img_path, _ in slides:
    s = prs.slides.add_slide(blank)
    # Fit image to full slide
    s.shapes.add_picture(img_path, Inches(0), Inches(0),
                         prs.slide_width, prs.slide_height)

out = '/home/xionghouyuan2/workplan_linear.pptx'
prs.save(out)
print(f'Created: {out}')
