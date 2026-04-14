#!/usr/bin/env python3
"""Daily English Study Plan - McKinsey Style PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import zipfile, os

# Colors
NAVY = RGBColor(0x05, 0x1C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_BLUE = RGBColor(0x00, 0x6B, 0xA6)
ACCENT_GREEN = RGBColor(0x00, 0x7A, 0x53)
ACCENT_ORANGE = RGBColor(0xD4, 0x6A, 0x00)

# Dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
LM = Inches(0.8)
CONTENT_W = Inches(11.733)

def add_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Remove style to prevent corruption
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, font_color=DARK_GRAY,
             bold=False, font_name='Arial', alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    from pptx.util import Pt as P
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = P(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.anchor = anchor
    return txBox

def add_hline(slide, left, top, length, color=BLACK, thickness=Pt(0.5)):
    add_rect(slide, left, top, length, thickness, color)

def add_oval(slide, left, top, text, size=Inches(0.5), bg=NAVY, fg=WHITE):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)
    # Add text
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = fg
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_action_title(slide, text):
    add_text(slide, LM, Inches(0.15), CONTENT_W, Inches(0.9), text,
             font_size=22, font_color=BLACK, bold=True, font_name='Georgia',
             anchor=MSO_ANCHOR.BOTTOM)
    add_hline(slide, LM, Inches(1.05), CONTENT_W, BLACK, Pt(0.5))

def add_source(slide, text):
    add_text(slide, LM, Inches(7.05), CONTENT_W, Inches(0.3), text,
             font_size=9, font_color=MED_GRAY)

def add_page_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=9, font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)

BL = prs.slide_layouts[6]

# ========== SLIDE 1: COVER ==========
s1 = prs.slides.add_slide(BL)
add_rect(s1, 0, 0, SW, Inches(0.05), NAVY)

# Title
add_text(s1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         'Daily English Study Plan',
         font_size=44, font_color=NAVY, bold=True, font_name='Georgia')

# Subtitle
add_text(s1, Inches(1), Inches(2.9), Inches(11), Inches(0.6),
         'A Structured 8-Week Program for English Mastery',
         font_size=24, font_color=DARK_GRAY)

# Date
add_text(s1, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
         '2026 | ClawBear 🐻',
         font_size=14, font_color=MED_GRAY)

# Decorative line
add_hline(s1, Inches(1), Inches(6.8), Inches(4), NAVY, Pt(2))

# ========== SLIDE 2: OVERVIEW ==========
s2 = prs.slides.add_slide(BL)
add_action_title(s2, 'Program Overview')

# Three pillars
pillars = [
    ('01', 'Vocabulary', 'Build a 5,000-word foundation with spaced repetition'),
    ('02', 'Listening', '30 minutes daily immersion with structured practice'),
    ('03', 'Speaking', 'Shadowing technique + conversation practice'),
]

pw = Inches(3.5)
pg = (CONTENT_W - pw * 3) / 2

for i, (num, title, desc) in enumerate(pillars):
    px = LM + (pw + pg) * i
    # Header bar
    add_rect(s2, px, Inches(1.5), pw, Inches(0.7), NAVY)
    add_oval(s2, px + Inches(0.15), Inches(1.58), num, size=Inches(0.5), bg=WHITE, fg=NAVY)
    add_text(s2, px + Inches(0.7), Inches(1.58), pw - Inches(0.9), Inches(0.5),
             title, font_size=16, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # Content box
    add_rect(s2, px, Inches(2.2), pw, Inches(3.5), BG_GRAY)
    add_text(s2, px + Inches(0.2), Inches(2.5), pw - Inches(0.4), Inches(2.8),
             desc, font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Key metrics
metrics = [('5,000+', 'Words'), ('56', 'Days'), ('2 hrs', 'Daily')]
mw = Inches(3.5)
mg = (CONTENT_W - mw * 3) / 2
for i, (val, label) in enumerate(metrics):
    mx = LM + (mw + mg) * i
    add_rect(s2, mx, Inches(6.0), mw, Inches(0.9), NAVY)
    add_text(s2, mx, Inches(6.0), mw, Inches(0.6),
             val, font_size=22, font_color=WHITE, bold=True, font_name='Georgia',
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s2, mx, Inches(6.55), mw, Inches(0.35),
             label, font_size=12, font_color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_source(s2, 'Daily English Study Plan')

# ========== SLIDE 3: DAILY SCHEDULE ==========
s3 = prs.slides.add_slide(BL)
add_action_title(s3, 'Daily Schedule')

# Timeline
schedule = [
    ('06:00-06:30', 'Morning Review', 'Flashcards + New words', ACCENT_BLUE),
    ('07:00-07:30', 'Listening Practice', 'Podcast or news', ACCENT_GREEN),
    ('12:00-12:15', 'Lunch Review', 'Anki sessions', ACCENT_ORANGE),
    ('19:00-20:00', 'Shadowing Session', 'Speak out loud', ACCENT_BLUE),
    ('20:00-20:30', 'Reading', 'Articles or books', ACCENT_GREEN),
    ('21:00-21:15', 'Evening Review', 'Day summary', ACCENT_ORANGE),
]

# Timeline line
add_hline(s3, LM + Inches(1.5), Inches(3.5), Inches(10), LINE_GRAY, Pt(2))

for i, (time, title, desc, color) in enumerate(schedule):
    tx = LM + (Inches(10) / (len(schedule) - 1)) * i
    ty = Inches(3.5)
    
    # Circle on timeline
    add_oval(s3, tx - Inches(0.15), ty - Inches(0.15), '', size=Inches(0.3), bg=color)
    
    # Time label above
    add_text(s3, tx - Inches(0.7), ty - Inches(1.0), Inches(1.4), Inches(0.5),
             time, font_size=11, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Content below
    box_h = Inches(1.8) if i % 2 == 0 else Inches(1.8)
    box_y = ty + Inches(0.5) if i % 2 == 0 else ty - Inches(2.3)
    
    add_rect(s3, tx - Inches(0.9), box_y, Inches(1.8), box_h, BG_GRAY)
    add_text(s3, tx - Inches(0.8), box_y + Inches(0.15), Inches(1.6), Inches(0.4),
             title, font_size=12, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s3, tx - Inches(0.8), box_y + Inches(0.55), Inches(1.6), Inches(0.9),
             desc, font_size=10, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_source(s3, 'Daily English Study Plan')
add_page_number(s3, 3, 8)

# ========== SLIDE 4: VOCABULARY ==========
s4 = prs.slides.add_slide(BL)
add_action_title(s4, 'Vocabulary Building')

# Method
add_rect(s4, LM, Inches(1.5), CONTENT_W, Inches(1.0), BG_GRAY)
add_text(s4, LM + Inches(0.3), Inches(1.6), Inches(2.0), Inches(0.5),
         'Method:', font_size=14, font_color=NAVY, bold=True)
add_text(s4, LM + Inches(2.3), Inches(1.6), CONTENT_W - Inches(2.6), Inches(0.8),
         'Spaced Repetition using Anki + Word Lists organized by topic',
         font_size=14, font_color=DARK_GRAY)

# Weekly topics
weeks = [
    ('Week 1-2', 'Daily Life & Routines'),
    ('Week 3-4', 'Work & Business'),
    ('Week 5-6', 'Technology & Innovation'),
    ('Week 7-8', 'Culture & Society'),
]

ww = Inches(2.7)
wg = (CONTENT_W - ww * 4) / 3

for i, (week, topic) in enumerate(weeks):
    wx = LM + (ww + wg) * i
    add_rect(s4, wx, Inches(2.8), ww, Inches(0.5), NAVY)
    add_text(s4, wx, Inches(2.8), ww, Inches(0.5),
             week, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(s4, wx, Inches(3.3), ww, Inches(2.5), BG_GRAY)
    add_text(s4, wx + Inches(0.15), Inches(3.5), ww - Inches(0.3), Inches(2.0),
             topic, font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Daily target
add_rect(s4, LM, Inches(6.0), CONTENT_W, Inches(0.8), NAVY)
add_text(s4, LM + Inches(0.3), Inches(6.0), Inches(3.0), Inches(0.8),
         'Daily Target:', font_size=14, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, LM + Inches(3.3), Inches(6.0), CONTENT_W - Inches(3.6), Inches(0.8),
         '20 new words + 50 review per day = 5,000+ words in 8 weeks',
         font_size=14, font_color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_source(s4, 'Daily English Study Plan')
add_page_number(s4, 4, 8)

# ========== SLIDE 5: LISTENING ==========
s5 = prs.slides.add_slide(BL)
add_action_title(s5, 'Listening Practice')

# Sources
sources = [
    ('Podcast', 'ESLPod, 6 Minute English', 'Easy'),
    ('News', 'BBC Learning English', 'Medium'),
    ('Talks', 'TED, Lex Fridman', 'Hard'),
    ('Series', 'Friends, The Office', 'Easy'),
]

sw = Inches(2.7)
sg = (CONTENT_W - sw * 4) / 3

for i, (category, examples, level) in enumerate(sources):
    sx = LM + (sw + sg) * i
    add_rect(s5, sx, Inches(1.5), sw, Inches(0.5), NAVY)
    add_text(s5, sx, Inches(1.5), sw, Inches(0.5),
             category, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(s5, sx, Inches(2.0), sw, Inches(2.5), BG_GRAY)
    add_text(s5, sx + Inches(0.15), Inches(2.2), sw - Inches(0.3), Inches(1.5),
             examples, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    # Level indicator
    level_color = ACCENT_GREEN if level == 'Easy' else (ACCENT_ORANGE if level == 'Medium' else ACCENT_BLUE)
    add_rect(s5, sx + Inches(0.85), Inches(4.0), sw - Inches(1.7), Inches(0.35), level_color)
    add_text(s5, sx + Inches(0.85), Inches(4.0), sw - Inches(1.7), Inches(0.35),
             level, font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Method box
add_rect(s5, LM, Inches(5.0), CONTENT_W, Inches(1.8), BG_GRAY)
add_text(s5, LM + Inches(0.3), Inches(5.1), CONTENT_W - Inches(0.6), Inches(0.4),
         'Practice Method:', font_size=14, font_color=NAVY, bold=True)
method_steps = [
    '1. Listen without subtitles first',
    '2. Listen with transcript',
    '3. Shadow the speaker (3x)',
    '4. Record yourself and compare'
]
add_text(s5, LM + Inches(0.3), Inches(5.5), CONTENT_W - Inches(0.6), Inches(1.2),
         '\n'.join(method_steps), font_size=12, font_color=DARK_GRAY)

add_source(s5, 'Daily English Study Plan')
add_page_number(s5, 5, 8)

# ========== SLIDE 6: SPEAKING ==========
s6 = prs.slides.add_slide(BL)
add_action_title(s6, 'Speaking Practice')

# Two methods
methods = [
    ('Shadowing', [
        'Choose a native speaker audio (podcast/TED)',
        'Listen one sentence at a time',
        'Pause and repeat immediately',
        'Copy the exact rhythm and intonation',
        'Do this 3 times before moving on'
    ]),
    ('Conversation', [
        'Find a language exchange partner (Tandem/HelloTalk)',
        'Or practice with AI (ChatGPT, Elsa Speak)',
        'Focus on fluency over accuracy',
        'Record and review your mistakes',
        'Aim for 30 minutes per session'
    ])
]

mw = Inches(5.5)
mg = Inches(0.733)

for i, (title, steps) in enumerate(methods):
    mx = LM + (mw + mg) * i
    add_rect(s6, mx, Inches(1.5), mw, Inches(0.5), NAVY)
    add_text(s6, mx, Inches(1.5), mw, Inches(0.5),
             title, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(s6, mx, Inches(2.0), mw, Inches(4.3), BG_GRAY)
    
    ty = Inches(2.2)
    for step in steps:
        add_oval(s6, mx + Inches(0.2), ty, '', size=Inches(0.25), bg=NAVY, fg=WHITE)
        add_text(s6, mx + Inches(0.55), ty, mw - Inches(0.8), Inches(0.6),
                 step, font_size=12, font_color=DARK_GRAY)
        ty += Inches(0.7)

# Daily goal
add_rect(s6, LM, Inches(6.5), CONTENT_W, Inches(0.6), NAVY)
add_text(s6, LM, Inches(6.5), CONTENT_W, Inches(0.6),
         'Daily Goal: Shadow for 20 min + Conversation for 10 min',
         font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_source(s6, 'Daily English Study Plan')
add_page_number(s6, 6, 8)

# ========== SLIDE 7: WEEKLY ROUTINE ==========
s7 = prs.slides.add_slide(BL)
add_action_title(s7, 'Weekly Focus')

days = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Weekend']
focuses = [
    (NAVY, 'Vocabulary\nHeavy'),
    (ACCENT_BLUE, 'Listening\nFocus'),
    (ACCENT_GREEN, 'Speaking\nFocus'),
    (NAVY, 'Vocabulary\nHeavy'),
    (ACCENT_ORANGE, 'Mixed\nPractice'),
    (BG_GRAY, 'Light Review\n+ Immersion'),
]

dw = Inches(1.8)
dg = (CONTENT_W - dw * 6) / 5

for i, (day, (color, focus)) in enumerate(zip(days, focuses)):
    dx = LM + (dw + dg) * i
    add_rect(s7, dx, Inches(1.5), dw, Inches(0.5), color if i < 5 else MED_GRAY)
    tc = WHITE if i < 5 else DARK_GRAY
    add_text(s7, dx, Inches(1.5), dw, Inches(0.5),
             day, font_size=12, font_color=tc, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(s7, dx, Inches(2.0), dw, Inches(2.5), BG_GRAY)
    add_text(s7, dx, Inches(2.2), dw, Inches(2.0),
             focus, font_size=13, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Weekly review box
add_rect(s7, LM, Inches(4.8), CONTENT_W, Inches(1.8), BG_GRAY)
add_text(s7, LM + Inches(0.3), Inches(4.9), Inches(2.0), Inches(0.4),
         'Weekend Review:', font_size=14, font_color=NAVY, bold=True)
weekend_items = [
    '• Review all new words from the week',
    '• Watch one English movie without subtitles',
    '• Write a short journal entry (100-200 words)',
    '• Prepare vocabulary list for next week'
]
add_text(s7, LM + Inches(0.3), Inches(5.3), CONTENT_W - Inches(0.6), Inches(1.2),
         '\n'.join(weekend_items), font_size=12, font_color=DARK_GRAY)

add_source(s7, 'Daily English Study Plan')
add_page_number(s7, 7, 8)

# ========== SLIDE 8: KEY TAKEAWAYS ==========
s8 = prs.slides.add_slide(BL)
add_action_title(s8, 'Key Takeaways')

takeaways = [
    ('Consistency', '30 minutes daily beats 3 hours once a week'),
    ('Input + Output', 'Listen AND speak — not just passive learning'),
    ('Spaced Repetition', 'Use Anki for efficient vocabulary retention'),
    ('Track Progress', 'Keep a log of what you learned each day'),
]

ty = Inches(1.5)
for i, (title, desc) in enumerate(takeaways):
    add_rect(s8, LM, ty, Inches(0.06), Inches(1.0), NAVY)
    add_oval(s8, LM + Inches(0.2), ty + Inches(0.1), str(i + 1), size=Inches(0.45), bg=NAVY)
    add_text(s8, LM + Inches(0.75), ty, Inches(3.0), Inches(0.4),
             title, font_size=16, font_color=NAVY, bold=True)
    add_text(s8, LM + Inches(0.75), ty + Inches(0.4), Inches(10), Inches(0.6),
             desc, font_size=14, font_color=DARK_GRAY)
    ty += Inches(1.2)

# Closing
add_rect(s8, 0, 0, SW, Inches(0.05), NAVY)
add_text(s8, Inches(1), Inches(5.5), Inches(11.3), Inches(1.0),
         'Start Today — Small Steps Lead to Big Results',
         font_size=28, font_color=NAVY, bold=True, font_name='Georgia',
         alignment=PP_ALIGN.CENTER)
add_hline(s8, Inches(5.5), Inches(6.8), Inches(2.3), NAVY, Pt(1.5))
add_text(s8, Inches(1), Inches(7.0), Inches(11.3), Inches(0.4),
         'Generated by ClawBear 🐻 | Daily English Study Plan',
         font_size=12, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Save
outpath = '/home/xionghouyuan2/.openclaw/workspace/daily-english-study-plan.pptx'
prs.save(outpath)

# Cleanup
tmppath = outpath + '.tmp'
with zipfile.ZipFile(outpath, 'r') as zin:
    with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith('.xml'):
                root = etree.fromstring(data)
                ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                for style in root.findall(f'.//{{{ns_p}}}style'):
                    style.getparent().remove(style)
            zout.writestr(item, data)
os.replace(tmppath, outpath)

print(f"Saved: {outpath}")
